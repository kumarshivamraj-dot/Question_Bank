from __future__ import annotations

from pathlib import Path

from study_pipeline.models import PageContent
from study_pipeline.text_utils import normalize_text


TEXT_PAGE_THRESHOLD = 80


def extract_document(path: Path, ocr_lang: str = "eng") -> list[PageContent]:
    suffix = path.suffix.lower()
    if suffix == ".pdf":
        return extract_pdf(path, ocr_lang=ocr_lang)
    if suffix == ".pptx":
        return extract_pptx(path)
    if suffix in {".txt", ".md"}:
        return [
            PageContent(
                document_path=str(path),
                document_name=path.name,
                page_number=1,
                text=path.read_text(encoding="utf-8"),
                extraction_method="plain_text",
            )
        ]
    raise ValueError(f"Unsupported file type: {path}")


def extract_pdf(path: Path, ocr_lang: str = "eng") -> list[PageContent]:
    import fitz

    from study_pipeline.ocr_fallback import ocr_image

    document = fitz.open(path)
    pages: list[PageContent] = []
    for index, page in enumerate(document, start=1):
        raw_text = normalize_text(page.get_text("text"))
        if len(raw_text) >= TEXT_PAGE_THRESHOLD:
            pages.append(
                PageContent(
                    document_path=str(path),
                    document_name=path.name,
                    page_number=index,
                    text=raw_text,
                    extraction_method="native_pdf_text",
                )
            )
            continue

        pixmap = page.get_pixmap(dpi=250, alpha=False)
        image = pixmap.pil_image()
        ocr_text = normalize_text(ocr_image(image, lang=ocr_lang))
        pages.append(
            PageContent(
                document_path=str(path),
                document_name=path.name,
                page_number=index,
                text=ocr_text,
                extraction_method="ocr_fallback",
            )
        )

    return pages


def extract_pptx(path: Path) -> list[PageContent]:
    from pptx import Presentation

    presentation = Presentation(path)
    pages: list[PageContent] = []
    for index, slide in enumerate(presentation.slides, start=1):
        text_runs: list[str] = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                text_runs.append(shape.text)
        pages.append(
            PageContent(
                document_path=str(path),
                document_name=path.name,
                page_number=index,
                text=normalize_text("\n".join(text_runs)),
                extraction_method="pptx_text",
            )
        )
    return pages
